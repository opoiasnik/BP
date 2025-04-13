from pptx import Presentation
from pptx.util import Inches

# Создание новой презентации
prs = Presentation()

# Добавляем пустой слайд (layout с индексом 5 обычно является пустым)
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)

# Определяем позицию и размер таблицы
left = Inches(0.5)
top = Inches(1.5)
width = Inches(9)
height = Inches(3)

# Количество строк: 1 заголовок + 2 строки с данными
rows = 3
cols = 4

# Добавляем таблицу на слайд
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Устанавливаем ширину столбцов (при необходимости можно настроить отдельно)
table.columns[0].width = Inches(1.5)  # Модель
table.columns[1].width = Inches(1)    # Оценка
table.columns[2].width = Inches(4)    # Текст
table.columns[3].width = Inches(2.5)  # Описание

# Заполняем заголовки
table.cell(0, 0).text = "Модель"
table.cell(0, 1).text = "Оценка"
table.cell(0, 2).text = "Текст"
table.cell(0, 3).text = "Описание"

# Данные для первого кандидата
table.cell(1, 0).text = "Mistral Small"
table.cell(1, 1).text = "9.0"
table.cell(1, 2).text = (
    "Nevolnosť môže byť spôsobená rôznymi príčinami, ako sú napríklad gastrointestinálne problémy, infekcie, alebo vedľajšie účinky liekov. "
    "Pre ľudí, ktorí hľadajú voľnopredajný liek na nevolnosť, sú dostupné niekoľko možností: 1. Dimedrol (Dramin) – Antihistaminikum; "
    "2. Bismut subsalicylát (Pepto-Bismol); 3. Ginger (Zázvor); 4. Meclizin (Bonine). Pred použitím lieku je dôležité konzultovať s lekárom."
)
table.cell(1, 3).text = "Evaluation based on required criteria."

# Данные для второго кандидата
table.cell(2, 0).text = "Mistral Large"
table.cell(2, 1).text = "8.0"
table.cell(2, 2).text = (
    "Pre nevolnosť sa dajú použiť niektoré voľne predávané lieky, ale dôležité je poradiť sa s lekárom, najmä ak má pacient 20 rokov. "
    "Medzi bežné voľne predávané lieky patria: 1. Dimenhydrinát; 2. Meclozín. Tieto lieky môžu spôsobiť spánkovosť a dávkovanie je nutné konzultovať."
)
table.cell(2, 3).text = "Evaluation based on required criteria."

# Сохраняем презентацию в файл
prs.save("evaluation_table.pptx")
