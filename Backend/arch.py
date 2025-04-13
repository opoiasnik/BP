from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR

# Vytvorenie novej prezentácie
prs = Presentation()
slide_layout = prs.slide_layouts[5]  # Prázdny slide
slide = prs.slides.add_slide(slide_layout)

# Definícia základných rozmerov a pozícií
left_margin = Inches(0.5)
top_margin = Inches(0.5)
block_width = Inches(3)
block_height = Inches(0.7)
vertical_gap = Inches(0.3)
horizontal_gap = Inches(0.5)

# Blok 1: Používateľský dotaz & Chat history
box1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left_margin, top_margin, block_width, block_height)
box1.text = "Používateľský dotaz\n& Chat history"

# Blok 2: ConversationalAgent (pod box1)
box2_top = top_margin + block_height + vertical_gap
box2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left_margin, box2_top, block_width, block_height)
box2.text = "ConversationalAgent"

# Blok 3: Klasifikácia dotazu (pod box2)
box3_top = box2_top + block_height + vertical_gap
box3 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left_margin, box3_top, block_width, block_height)
box3.text = "Klasifikácia dotazu"

# Vetvenie: Pozície pre dve vetvy
branch_top = box3_top + block_height + vertical_gap

# Ľavá vetva ("Vyhladavanie")
left_branch_left = left_margin - Inches(0.2)
box4A = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left_branch_left, branch_top, block_width, block_height)
box4A.text = "ElasticsearchStore\nvyhľadávanie"

box5A_top = branch_top + block_height + vertical_gap
box5A = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left_branch_left, box5A_top, block_width, block_height)
box5A.text = "Generovanie\ndynamického promptu"

box6A_top = box5A_top + block_height + vertical_gap
box6A = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left_branch_left, box6A_top, block_width, block_height)
box6A.text = "Generovanie\nodpovede"

box7A_top = box6A_top + block_height + vertical_gap
box7A = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left_branch_left, box7A_top, block_width, block_height)
box7A.text = "Finalizácia\nodpovede"

# Pravá vetva ("Upresnenie")
right_branch_left = left_margin + block_width + horizontal_gap
box4B = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, right_branch_left, branch_top, block_width, block_height)
box4B.text = "Kombinovanie\ndotazov"

box5B_top = branch_top + block_height + vertical_gap
box5B = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, right_branch_left, box5B_top, block_width, block_height)
box5B.text = "ElasticsearchStore\nvyhľadávanie"

box6B_top = box5B_top + block_height + vertical_gap
box6B = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, right_branch_left, box6B_top, block_width, block_height)
box6B.text = "Generovanie\ndynamického promptu"

box7B_top = box6B_top + block_height + vertical_gap
box7B = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, right_branch_left, box7B_top, block_width, block_height)
box7B.text = "Generovanie\nodpovedí (2 modely)"

box8B_top = box7B_top + block_height + vertical_gap
box8B = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, right_branch_left, box8B_top, block_width, block_height)
box8B.text = "Validácia a\nhodnotenie odpovedí"

box9B_top = box8B_top + block_height + vertical_gap
box9B = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, right_branch_left, box9B_top, block_width, block_height)
box9B.text = "Finalizácia\nodpovede"

# Finálny blok: Výstup (zlúčenie vetiev)
final_box_top = max(box7A_top, box9B_top) + block_height + vertical_gap
final_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left_margin, final_box_top, block_width, block_height)
final_box.text = "Výstup"

# Funkcia na pridanie šípok medzi blokmi
def add_connector(slide, start_shape, end_shape):
    start_x = start_shape.left + start_shape.width / 2
    start_y = start_shape.top + start_shape.height
    end_x = end_shape.left + end_shape.width / 2
    end_y = end_shape.top
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y)
    # Aktuálna verzia python-pptx nepodporuje nastavenie šípky, preto tento riadok odstraňte alebo zakomentujte:
    # connector.line.end_arrowhead.style = 1
    return connector

# Prepojenie blokov
add_connector(slide, box1, box2)
add_connector(slide, box2, box3)

# Vetvenie z Box3 do oboch vetiev
mid_point = box3.left + box3.width / 2
branch_mid_y = box3.top + box3.height + vertical_gap/2
# Do ľavej vetvy:
connector_left = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, mid_point, box3.top + box3.height, left_branch_left + block_width/2, branch_mid_y)
# Do pravej vetvy:
connector_right = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, mid_point, box3.top + box3.height, right_branch_left + block_width/2, branch_mid_y)

# Prepojenie blokov v ľavej vetve
add_connector(slide, box4A, box5A)
add_connector(slide, box5A, box6A)
add_connector(slide, box6A, box7A)

# Prepojenie blokov v pravej vetve
add_connector(slide, box4B, box5B)
add_connector(slide, box5B, box6B)
add_connector(slide, box6B, box7B)
add_connector(slide, box7B, box8B)
add_connector(slide, box8B, box9B)

# Spojenie oboch vetiev s finálnym blokom "Výstup"
add_connector(slide, box7A, final_box)
add_connector(slide, box9B, final_box)

# Uloženie prezentácie
prs.save("architecture_diagram.pptx")
